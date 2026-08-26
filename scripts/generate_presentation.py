#!/usr/bin/env python3
"""
Generates a presentation deck from a branding profile.

Usage:
    python3 scripts/generate_presentation.py               # generic profile
    python3 scripts/generate_presentation.py \\
        --profile <path-to-profile.py>                     # custom profile

Requires: pip install python-pptx
"""

import argparse
import importlib.util
from pathlib import Path


def _load_profile(path):
    """Load BRAND dict from a profile file, or fall back to profile_generic."""
    if path:
        spec = importlib.util.spec_from_file_location("profile", path)
    else:
        here = Path(__file__).parent
        spec = importlib.util.spec_from_file_location(
            "profile_generic",
            here / "presentation" / "profile_generic.py",
        )
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod.BRAND


_ap = argparse.ArgumentParser(description="Generate a presentation deck.")
_ap.add_argument("--profile", default=None,
                 help="Path to a profile .py exporting BRAND dict "
                      "(default: scripts/presentation/profile_generic.py).")
BRAND = _load_profile(_ap.parse_args().profile)

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Brand palette — sourced from profile ─────────────────────
PRIMARY       = RGBColor(*BRAND["primary"])
PRIMARY_DARK  = RGBColor(*BRAND["primary_dark"])
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x1A)
MID       = RGBColor(0x55, 0x55, 0x55)
LIGHT     = RGBColor(0xF2, 0xF2, 0xF2)
GREEN     = RGBColor(0x00, 0x7A, 0x3D)
AMBER     = RGBColor(0xFF, 0x8C, 0x00)
BLUE      = RGBColor(0x00, 0x4E, 0x9A)
BORDER    = RGBColor(0xCC, 0xCC, 0xCC)
STRIPE    = RGBColor(0xF8, 0xF8, 0xF8)
PRIMARY_TINT  = RGBColor(*BRAND["primary_tint"])


_NOTE_SUBS = BRAND.get("note_subs", [])


def _b(s):
    """Apply brand substitutions to speaker-notes strings at runtime.
    Substitution pairs come from BRAND["note_subs"] in the active profile.
    For the generic profile, this is a no-op.
    For the internal profile, restores organisation-specific copy in notes.
    """
    for generic, branded in _NOTE_SUBS:
        s = s.replace(generic, branded)
    return s

# ── Canvas: 16:9 widescreen (13.33" × 7.5") ──────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


def R(sl, x, y, w, h, fill, *, border=False, bc=None, bw=0.75):
    """Add a solid rectangle."""
    shp = sl.shapes.add_shape(1, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border and bc:
        shp.line.color.rgb = bc
        shp.line.width = Pt(bw)
    else:
        shp.line.fill.background()
    return shp


def T(sl, text, x, y, w, h, *,
      size=14, color=DARK, bold=False, italic=False,
      align=PP_ALIGN.LEFT, font="Calibri", wrap=True):
    """Add a text box with a single run."""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return tb


def header(sl, title):
    """Draw standard red header bar + title. Returns content_top y-position."""
    R(sl, Inches(0), Inches(0), W, H, LIGHT)
    R(sl, Inches(0), Inches(0), W, Inches(1.1), PRIMARY)
    T(sl, title,
      Inches(0.45), Inches(0.17), Inches(12.5), Inches(0.9),
      size=25, color=WHITE, bold=True)
    R(sl, Inches(0), Inches(1.1), W, Inches(0.05), PRIMARY_DARK)
    return Inches(1.25)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
R(sl, Inches(0), Inches(0), W, H, PRIMARY)
R(sl, Inches(0), H - Inches(2.15), W, Inches(2.15), PRIMARY_DARK)

T(sl, BRAND["org_upper"],
  Inches(0.55), Inches(0.30), Inches(6), Inches(0.5),
  size=13, color=WHITE, bold=True, font="Calibri Light")
R(sl, Inches(0.55), Inches(0.82), Inches(11.8), Inches(0.022), WHITE)

T(sl, f"Web Intelligence\nfor {BRAND['platform']}",
  Inches(0.55), Inches(1.05), Inches(9), Inches(3.1),
  size=50, color=WHITE, bold=True, font="Calibri Light")

T(sl, "Self-hosted · POPIA-compliant · MCP-native\n"
       "Web search & content extraction for enterprise AI agents — "
       f"on infrastructure {BRAND['org']} already owns.",
  Inches(0.55), Inches(4.2), Inches(9.6), Inches(1.05),
  size=17, color=RGBColor(0xFF, 0xCC, 0xCC), font="Calibri Light")

for i, (val, lbl) in enumerate([
    ("$0",   "marginal infra on OpenShift"),
    ("MIT",  "open-source licence"),
    ("POPIA", "compliant by design"),
]):
    bx = Inches(0.55) + i * Inches(4.22)
    T(sl, val, bx, H - Inches(1.92), Inches(4.0), Inches(0.72),
      size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, lbl, bx, H - Inches(1.2), Inches(4.0), Inches(0.38),
      size=11, color=RGBColor(0xFF, 0xCC, 0xCC), align=PP_ALIGN.CENTER)

T(sl, BRAND["prepared_by"],
  Inches(0.55), H - Inches(0.32), Inches(7), Inches(0.28),
  size=9, color=RGBColor(0xFF, 0xAA, 0xAA), font="Calibri Light")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "The Opportunity: AI Agents That Can See the World")

T(sl, f"{BRAND['platform']} agents are powerful reasoners — but without live web access and "
       "internal knowledge, every answer is bounded by training data that is already months old.",
  Inches(0.5), ct, Inches(12.4), Inches(0.55),
  size=14, color=MID)

cards = [
    ("Web Fetch",
     ["Retrieve any URL as clean, LLM-ready Markdown.",
      "Playwright handles JavaScript-rendered pages correctly.",
      "Tables, code blocks, and headings preserved — minimal information loss.",
      "Configurable resource blocking reduces page weight ~60%."],
     PRIMARY),
    ("Web Search",
     ["Resolve natural-language queries to ranked URLs.",
      "Supports SearXNG (self-hosted), Brave API, DuckDuckGo.",
      "Multi-provider fan-out with Reciprocal Rank Fusion merge.",
      "Per-engine POPIA profile: choose the right data boundary."],
     BLUE),
    ("Enterprise Knowledge Index",
     ["Crawl SharePoint, Confluence, internal portals.",
      "Cache LLM-ready Markdown from your own document corpus.",
      "Serve results in milliseconds — no live API calls.",
      "Respects per-user access controls (Phase 2)."],
     AMBER),
]

cw = Inches(3.98)
gap = Inches(0.2)
sx = Inches(0.5)
cy = ct + Inches(0.7)

for i, (title, bullets, c) in enumerate(cards):
    cx = sx + i * (cw + gap)
    shp = sl.shapes.add_shape(1, cx, cy, cw, Inches(4.0))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = BORDER; shp.line.width = Pt(0.5)
    R(sl, cx, cy, cw, Inches(0.07), c)
    T(sl, title, cx + Inches(0.18), cy + Inches(0.14), cw - Inches(0.36), Inches(0.42),
      size=15, color=c, bold=True)
    for bi, b in enumerate(bullets):
        by = cy + Inches(0.65) + bi * Inches(0.78)
        R(sl, cx + Inches(0.18), by + Inches(0.18), Inches(0.09), Inches(0.09), c)
        T(sl, b, cx + Inches(0.35), by, cw - Inches(0.53), Inches(0.72),
          size=12, color=MID)
    if i == 2:
        R(sl, cx + cw - Inches(1.5), cy + Inches(3.55), Inches(1.35), Inches(0.3), AMBER)
        T(sl, "Roadmap →", cx + cw - Inches(1.5), cy + Inches(3.55), Inches(1.35), Inches(0.3),
          size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHY NOT EXTERNAL VENDORS
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, f"Why External Vendors Are Not the Answer for {BRAND['org']}")

issues = [
    (PRIMARY,   "POPIA Section 72 — Cross-Border Transfer",
     "Every query sent to Tavily (US) or Firecrawl (US) is a cross-border transfer of personal information. "
     "For a regulated SA operator, this requires a binding Section 72 justification that may be "
     "unavailable for all data categories. Self-hosted software on SA infrastructure eliminates this risk entirely."),
    (AMBER, "Unpredictable Cost at Scale",
     "Tavily Growth = $0.005/query. At 10,000 queries/day: ~$18,250/year (≈R330,000) before any extract workload. "
     "Firecrawl Standard → Growth tier jump: $83/mo to $333/mo — a 4× increase overnight. "
     "A workload crossing a pricing tier triggers an immediate jump to the next tier's flat rate with no warning."),
    (AMBER, "Vendor ToS and Availability Risk",
     "Free search endpoints (DuckDuckGo, public SearXNG) provide no SLA and carry documented ToS exposure. "
     "Any vendor pricing or ToS change mid-contract requires an immediate architectural response. "
     "Self-hosted deployments are fully insulated from vendor decisions."),
    (BLUE,  "No Governance Pack Included",
     "Neither Tavily nor Firecrawl ships a POPIA assessment, data flow inventory, threat model, or runbook. "
     f"{BRAND['org_poss']} security and privacy review process requires all of these artefacts before production approval. "
     "This tool's governance pack is complete and already in the repository — ready for review today."),
]

rh = Inches(1.2)
gap = Inches(0.07)
for i, (color, title, body) in enumerate(issues):
    ry = ct + i * (rh + gap)
    R(sl, Inches(0.5),  ry, Inches(0.14), rh, color)
    R(sl, Inches(0.66), ry, Inches(12.12), rh, WHITE)
    T(sl, title, Inches(0.86), ry + Inches(0.1), Inches(11.6), Inches(0.38),
      size=13, color=DARK, bold=True)
    T(sl, body,  Inches(0.86), ry + Inches(0.5), Inches(11.6), Inches(0.65),
      size=11, color=MID)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT IS THIS TOOL
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "Introducing: markdown-for-agents-mcp")

T(sl, f"A lightweight MCP server exposing web fetch, web search, and (Phase 2) enterprise "
       f"knowledge index as tools your {BRAND['platform']} agents call via the Model Context Protocol.",
  Inches(0.5), ct, Inches(12.4), Inches(0.55),
  size=14, color=MID)

lx = Inches(0.5)
lw_label = Inches(1.55)
lw_body  = Inches(4.3)
cy2 = ct + Inches(0.72)

facts = [
    ("Protocol",     f"MCP (Model Context Protocol) — an emerging standard for agent tool calls, adopted by Anthropic, OpenAI, and major LLM frameworks. {BRAND['platform']} connects via SSE or stdio."),
    ("Dependencies", "7 runtime dependencies. No managed Redis, no RabbitMQ, no PostgreSQL required. Firecrawl self-hosted: 6–7 services. Simpler = fewer failure modes."),
    ("Licence",      f"MIT. No AGPL copyleft, no CLA. Can be modified and embedded in {BRAND['org_poss']} agent platform without restriction. Firecrawl is AGPL-3.0."),
    ("Renderer",     "Playwright (Chromium) for JavaScript-rendered pages. Handles modern SPAs, SharePoint modern pages, and Confluence macros correctly."),
    ("Deployment",   "Docker Compose → ECS Fargate (Mode F, " + BRAND["platform"] + "'s existing runtime) → OpenShift (Mode G, $0 marginal infra on existing cluster)."),
]

for i, (label, body) in enumerate(facts):
    fy = cy2 + i * Inches(1.02)
    R(sl, lx, fy, lw_label, Inches(0.84), PRIMARY)
    T(sl, label, lx + Inches(0.08), fy + Inches(0.2), lw_label - Inches(0.16), Inches(0.46),
      size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    R(sl, lx + lw_label + Inches(0.06), fy, lw_body, Inches(0.84), WHITE,
      border=True, bc=BORDER, bw=0.4)
    T(sl, body, lx + lw_label + Inches(0.18), fy + Inches(0.07), lw_body - Inches(0.25), Inches(0.72),
      size=11, color=MID)

# Right column: stat highlights
rx = Inches(7.0)
rw = Inches(5.8)
highlights = [
    ("7",       "npm runtime dependencies\nvs 6–7 deployed services for\nFirecrawl self-hosted",  PRIMARY),
    ("MIT",     "open-source licence\nno AGPL copyleft restrictions",       GREEN),
    ("<5 min",  f"to connect {BRAND['platform']}\nvia MCP",                            BLUE),
    ("1,237",   "tests passing in CI\nproduction-grade quality",            PRIMARY),
]
hw = Inches(2.55)
hh = Inches(1.35)
for i, (val, lbl, c) in enumerate(highlights):
    hx = rx + (i % 2) * (hw + Inches(0.35))
    hy = cy2 + (i // 2) * (hh + Inches(0.15))
    shp = sl.shapes.add_shape(1, hx, hy, hw, hh)
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = BORDER; shp.line.width = Pt(0.5)
    R(sl, hx, hy, hw, Inches(0.06), c)
    T(sl, val, hx, hy + Inches(0.1), hw, Inches(0.66),
      size=30, color=c, bold=True, align=PP_ALIGN.CENTER)
    T(sl, lbl, hx, hy + Inches(0.78), hw, Inches(0.52),
      size=11, color=MID, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "How It Works: The Rendering Pipeline")

T(sl, f"A query from a {BRAND['platform']} agent becomes LLM-ready Markdown in five deterministic steps — "
       "with a shared URL cache that avoids redundant fetches across agent sessions.",
  Inches(0.5), ct, Inches(12.4), Inches(0.45),
  size=14, color=MID)

steps = [
    ("1  Query",   "Agent calls\nfetch() or\nsearch()"),
    ("2  Resolve", "Provider chosen\n(cache / SearXNG /\nBrave / Graph)"),
    ("3  Render",  "Render tier chosen:\nHTTP fast-path;\nescalates to\nPlaywright if JS-heavy."),
    ("4  Extract", "markdown-for-agents\nlibrary → structured,\nboilerplate-free\nMarkdown."),
    ("5  Return",  "Full Markdown\nreturned with\nsource URL &\nmetadata."),
]

sw = Inches(2.12)
sh = Inches(2.65)
gap = Inches(0.22)
sy = ct + Inches(0.62)
sx = Inches(0.5)

for i, (title, body) in enumerate(steps):
    bx = sx + i * (sw + gap)
    if i < len(steps) - 1:
        T(sl, "→", bx + sw + Inches(0.04), sy + sh / 2 - Inches(0.25), gap, Inches(0.45),
          size=16, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    shp = sl.shapes.add_shape(1, bx, sy, sw, sh)
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = BORDER; shp.line.width = Pt(0.5)
    R(sl, bx, sy, sw, Inches(0.07), PRIMARY)
    T(sl, title, bx, sy + Inches(0.08), sw, Inches(0.5),
      size=13, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    T(sl, body, bx + Inches(0.1), sy + Inches(0.65), sw - Inches(0.2), Inches(1.85),
      size=11, color=MID, align=PP_ALIGN.CENTER)

# Feature row
feats = [
    ("Shared URL cache",   "Avoids re-fetching identical URLs within the TTL window — reduces egress and latency across all agent sessions."),
    ("Resource blocking",  "Chromium blocks ads, fonts, and tracking scripts by default. Reduces median page weight by ~60%."),
    ("Markdown fidelity",  "Tables, code blocks, headings, and lists all preserved. Structure faithfully converted — minimal information loss vs raw HTML."),
    ("Provider fan-out",   "Multiple search providers queried in parallel. Reciprocal Rank Fusion merges ranked results into one list."),
]
fy = sy + sh + Inches(0.28)
fw = Inches(2.97)
for i, (title, body) in enumerate(feats):
    fx = Inches(0.5) + i * (fw + Inches(0.17))
    R(sl, fx, fy, fw, Inches(1.2), LIGHT, border=True, bc=BORDER, bw=0.4)
    R(sl, fx, fy, Inches(0.08), Inches(1.2), PRIMARY)
    T(sl, title, fx + Inches(0.18), fy + Inches(0.1), fw - Inches(0.26), Inches(0.34),
      size=11, color=DARK, bold=True)
    T(sl, body, fx + Inches(0.18), fy + Inches(0.46), fw - Inches(0.26), Inches(0.68),
      size=10, color=MID)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — POPIA & DATA SOVEREIGNTY
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "POPIA & Data Sovereignty: Built In, Not Bolted On")

lx = Inches(0.5)
lw = Inches(6.1)
points = [
    ("Section 72 — Cross-Border Transfer Eliminated",
     f"Self-hosted on {BRAND['org_poss']} SA infrastructure means no personal information leaves South African borders. "
     "No Section 72 justification is required — this is the cleanest possible POPIA posture, "
     "better than any cloud-managed service including AWS af-south-1 or Azure South Africa North."),
    ("Section 19 — Full Auditability",
     "Every data flow is documented in DATA_FLOW.md. Every component is open-source and inspectable. "
     "No black-box SaaS. Security and privacy reviewers can verify the complete stack without "
     "depending on a vendor's word."),
    ("Three Configurable Engine Profiles",
     "'clean': zero external calls (target URLs only). 'balanced': SearXNG self-hosted. "
     "'external': Brave API with DPA. Each profile has a documented data flow and a POPIA risk rating. "
     "Operators select the profile that matches their legal advice."),
]
py = ct + Inches(0.1)
for title, body in points:
    R(sl, lx, py, Inches(0.1), Inches(1.5), PRIMARY)
    T(sl, title, lx + Inches(0.22), py + Inches(0.08), lw - Inches(0.32), Inches(0.4),
      size=13, color=DARK, bold=True)
    T(sl, body, lx + Inches(0.22), py + Inches(0.5), lw - Inches(0.32), Inches(0.94),
      size=11, color=MID)
    py += Inches(1.65)

# Right: governance table
rx = Inches(7.1)
rw = Inches(5.75)
T(sl, "Governance Pack — Included in the Repository",
  rx, ct, rw, Inches(0.42),
  size=13, color=DARK, bold=True)

gov_items = [
    ("POPIA_ASSESSMENT.md",          "Full POPIA impact assessment"),
    ("DATA_FLOW.md",                 "Data flow inventory per engine profile"),
    ("THREAT_MODEL.md",              "Threat model with explicit ceilings"),
    ("TERMS_OF_SERVICE.md",          "ToS analysis per search provider"),
    ("SECURITY_FINDINGS_REGISTER.md","Open security findings tracker"),
    ("ENTERPRISE_READINESS.md",      "Honest readiness assessment"),
    ("RUNBOOK.md",                   "Operations runbook"),
    ("SLO.md",                       "SLO template (measured at deployment)"),
    ("COST_ANALYSIS.md",             "Full TCO model & build/buy analysis"),
    ("PRODUCTION_AUTHORISATION.md",  "Authorisation gate (status: pre-production — pending deployment)"),
]

tbl_top = ct + Inches(0.52)
tbl_h   = Inches(0.43)
for i, (doc, desc) in enumerate(gov_items):
    ty = tbl_top + i * (tbl_h + Inches(0.03))
    bg = LIGHT if i % 2 == 0 else WHITE
    R(sl, rx, ty, rw, tbl_h, bg)
    R(sl, rx, ty, Inches(0.06), tbl_h, PRIMARY)
    T(sl, doc,  rx + Inches(0.12), ty + Inches(0.07), Inches(2.5), Inches(0.3),
      size=9, color=PRIMARY, bold=True, font="Courier New")
    T(sl, desc, rx + Inches(2.65), ty + Inches(0.07), Inches(2.9), Inches(0.3),
      size=10, color=MID)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TOTAL COST OF OWNERSHIP
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "Total Cost of Ownership vs Alternatives")

T(sl, "At 100,000 pages/month. All figures from the reproducible model in scripts/cost-model.py. "
       "Engineering at R1.5M/yr fully-loaded SA senior rate. Mode G = existing OpenShift cluster ($0 new infra).",
  Inches(0.5), ct, Inches(12.4), Inches(0.45),
  size=11, color=MID)

col_w = [Inches(3.05), Inches(1.95), Inches(2.05), Inches(1.95), Inches(3.3)]
headers = ["Option", "Infrastructure", "Engineering", "Total / month", "POPIA posture"]
rows_d = [
    ["Tavily Growth  (100k queries)",      "$500/mo",    "Minimal",           "~$500/mo",         "⚠  Risk — US SaaS"],
    ["Firecrawl Standard  (100k pages)",   "$83/mo",     "Minimal",           "~$83/mo",          "⚠  Risk — US SaaS"],
    ["AWS Kendra  (af-south-1, 200k docs)","$1,981/mo",  "Moderate",          ">$2,000/mo",       "✓  Acceptable — SA region"],
    ["Glean Enterprise  (500 users)",      "$0 (your cloud)", "Managed SaaS", "~$25,000/mo",      "✗  Limited — managed access"],
    ["This tool — Mode F  (ECS Fargate)",  "$251/mo",    "$667/mo (0.1 FTE)", "~$918/mo",         "✓  Clean — self-hosted SA"],
    ["This tool — Mode G  (OpenShift)",    "$0",         "$100–$527/mo",      "~$100–$527/mo",    "✓✓ Best — self-hosted SA"],
]

tbl = sl.shapes.add_table(
    len(rows_d) + 1, len(headers),
    Inches(0.5), ct + Inches(0.52),
    sum(col_w), Inches(3.55)
).table

for ci, (w, h) in enumerate(zip(col_w, headers)):
    tbl.columns[ci].width = w
    cell = tbl.cell(0, ci)
    cell.text = h
    cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
    p = cell.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.color.rgb = WHITE; run.font.bold = True
    run.font.size = Pt(11); run.font.name = "Calibri"

hl = {4, 5}
for ri, row in enumerate(rows_d):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri + 1, ci)
        cell.text = val
        is_hl = ri in hl
        bg = PRIMARY_TINT if is_hl else (STRIPE if ri % 2 == 0 else WHITE)
        cell.fill.solid(); cell.fill.fore_color.rgb = bg
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
        run = p.runs[0]
        run.font.size = Pt(10); run.font.name = "Calibri"
        run.font.color.rgb = PRIMARY if (is_hl and ci in (0, 3)) else DARK
        run.font.bold = is_hl and ci == 0

# Callout below table
cy3 = ct + Inches(4.15)
R(sl, Inches(0.5), cy3, Inches(12.33), Inches(0.88), PRIMARY_TINT)
R(sl, Inches(0.5), cy3, Inches(0.1), Inches(0.88), PRIMARY)
T(sl, "Mode G ongoing: $100/mo. Firecrawl Standard: $83/mo covers up to 100k pages/month — "
       "Mode G is $17/mo more at this volume. Exceed 100k pages and Firecrawl jumps to Growth ($333/mo), "
       "making Mode G the cheaper option. Vs Glean Enterprise (500 users): ~250× cheaper at any volume.",
  Inches(0.7), cy3 + Inches(0.08), Inches(12.0), Inches(0.72),
  size=11, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DEPLOYMENT ON EXISTING INFRASTRUCTURE
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, f"Runs on Infrastructure {BRAND['org']} Already Owns")

modes = [
    ("Mode F\nECS Fargate",
     BRAND["platform"] + "'s Existing Runtime",
     [BRAND["platform"] + " already runs here — same pipeline",
      "Fargate Spot reduces compute cost by ~70%",
      "Right-sized: ~$251/mo infrastructure",
      "Scales to HPA max — set a spend cap",
      "af-south-1 (POPIA data residency)"],
     BLUE, "Existing"),
    ("Mode G\nOpenShift",
     "Recommended Path",
     ["$0 marginal infrastructure cost",
      "Runs in existing OCP cluster",
      "OpenShift Routes + SecurityContextConstraints",
      "KEDA for autoscaling — no additional cost",
      "Best POPIA posture: fully on-cluster"],
     GREEN, "Recommended"),
    (BRAND["llm_platform"] + "\nIn-house LLM",
     "Complementary Capability",
     [BRAND["llm_platform"] + " handles LLM inference",
      "This tool handles web retrieval",
      "MCP bridges both cleanly",
      "No capability overlap",
      "Together: fully self-contained AI platform"],
     AMBER, "Complementary"),
]

mw = Inches(3.85)
gap = Inches(0.35)
sx = Inches(0.5)
my = ct + Inches(0.25)

for i, (title, subtitle, bullets, c, badge) in enumerate(modes):
    mx = sx + i * (mw + gap)
    shp = sl.shapes.add_shape(1, mx, my, mw, Inches(4.95))
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = c; shp.line.width = Pt(1.5)
    R(sl, mx, my, mw, Inches(0.08), c)
    R(sl, mx + mw - Inches(1.55), my + Inches(0.14), Inches(1.4), Inches(0.28), c)
    T(sl, badge, mx + mw - Inches(1.55), my + Inches(0.14), Inches(1.4), Inches(0.28),
      size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, title, mx + Inches(0.18), my + Inches(0.1), mw - Inches(1.65), Inches(0.72),
      size=15, color=c, bold=True)
    T(sl, subtitle, mx + Inches(0.18), my + Inches(0.84), mw - Inches(0.36), Inches(0.32),
      size=11, color=MID, italic=True)
    R(sl, mx + Inches(0.18), my + Inches(1.22), mw - Inches(0.36), Inches(0.02), c)
    for bi, b in enumerate(bullets):
        by = my + Inches(1.35) + bi * Inches(0.68)
        R(sl, mx + Inches(0.2), by + Inches(0.18), Inches(0.09), Inches(0.09), c)
        T(sl, b, mx + Inches(0.4), by, mw - Inches(0.58), Inches(0.62),
          size=12, color=MID)

R(sl, Inches(0.5), ct + Inches(5.38), Inches(12.33), Inches(0.62), LIGHT)
R(sl, Inches(0.5), ct + Inches(5.38), Inches(0.1), Inches(0.62), GREEN)
T(sl, f"Mode G recommendation: $0 marginal infra, deployable in days, runs alongside {BRAND['platform']} "
       f"on the existing OpenShift cluster. No new procurement. No new approval cycle for {BRAND['org']}.",
  Inches(0.7), ct + Inches(5.45), Inches(12.0), Inches(0.48),
  size=11, color=MID)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PLATFORM INTEGRATION
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, f"Integration with {BRAND['platform']} via MCP")

T(sl, f"The Model Context Protocol is an emerging standard for connecting AI agents to external tools. "
       f"{BRAND['platform']} connects to the MCP server via standard configuration — no bespoke API integration code required.",
  Inches(0.5), ct, Inches(12.4), Inches(0.5),
  size=14, color=MID)

stack = [
    (BRAND["platform"] + " Agent",  "LLM reasoning + task execution via " + BRAND["llm_platform"],  BLUE),
    ("MCP Tool Call",            "fetch()  ·  search()  ·  (Phase 2: search_knowledge())",   PRIMARY),
    ("markdown-for-agents-mcp",  "MCP server — tool dispatch · rendering · index queries",  DARK),
    ("Source Connectors",        "SearXNG · Brave API  (Phase 2: SharePoint Graph API · Confluence)",  MID),
    ("Infrastructure",           "OpenShift Mode G  (existing cluster, $0 marginal cost)",  GREEN),
]

bw = Inches(11.3)
bh = Inches(0.72)
bx = Inches(1.0)
by = ct + Inches(0.65)

for i, (title, subtitle, c) in enumerate(stack):
    R(sl, bx, by, bw, bh, c)
    T(sl, title, bx + Inches(0.2), by + Inches(0.1), Inches(3.6), Inches(0.52),
      size=14, color=WHITE, bold=True)
    T(sl, subtitle, bx + Inches(3.9), by + Inches(0.13), Inches(7.2), Inches(0.45),
      size=12, color=WHITE if c not in (LIGHT, MID) else WHITE)
    if i < len(stack) - 1:
        T(sl, "↕", bx + bw / 2 - Inches(0.2), by + bh, Inches(0.4), Inches(0.22),
          size=12, color=MID, align=PP_ALIGN.CENTER)
    by += bh + Inches(0.22)

# Side labels
labels = [
    (Inches(0.08), ct + Inches(0.65),  Inches(0.85), Inches(0.72), "AGENT",   BLUE),
    (Inches(0.08), ct + Inches(4.41),  Inches(0.85), Inches(0.72), "INFRA",   GREEN),
]
for lx, ly, lw2, lh, lbl, c in labels:
    R(sl, lx, ly, lw2, lh, c)
    T(sl, lbl, lx, ly + Inches(0.22), lw2, Inches(0.3),
      size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DESIGNED FOR AGENTIC SYSTEMS
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "Designed for Agentic Systems — Not Individual Sessions")

T(sl, "Claude Code's built-in web tools are designed for one developer in one session. "
       f"{BRAND['platform']} runs many concurrent agents. The difference is architectural.",
  Inches(0.5), ct, Inches(12.4), Inches(0.45),
  size=14, color=MID)

# Left column: the problem with per-agent tools
lx = Inches(0.5)
lw = Inches(5.85)
ly = ct + Inches(0.58)
R(sl, lx, ly, lw, Inches(0.36), MID)
T(sl, "Per-Agent Built-in Tools  (Claude Code default)",
  lx + Inches(0.15), ly + Inches(0.07), lw - Inches(0.3), Inches(0.25),
  size=11, color=WHITE, bold=True)

problems = [
    ("No shared state", "Each agent fetches the same URL independently. 10 concurrent agents fetching the same page = 10 separate network fetches. No shared URL cache between sessions or agents."),
    ("No render intelligence", "Every call is stateless. There is no equivalent to the three-tier render ladder or domain-level tier memoization — each fetch is a fresh HTTP request with no learned context."),
    ("Uncontrolled egress", "Each agent's web calls go out independently. No single network chokepoint for POPIA audit, rate limiting, or provider switching."),
    ("Output variability", "Raw or minimally processed HTML enters the agent's context directly. Quality depends entirely on the site's markup — inconsistent across agents."),
]
py = ly + Inches(0.45)
for title, body in problems:
    R(sl, lx, py, lw, Inches(0.88), WHITE)
    R(sl, lx, py, Inches(0.08), Inches(0.88), AMBER)
    T(sl, title, lx + Inches(0.2), py + Inches(0.06), lw - Inches(0.28), Inches(0.3),
      size=12, color=DARK, bold=True)
    T(sl, body,  lx + Inches(0.2), py + Inches(0.38), lw - Inches(0.28), Inches(0.46),
      size=10, color=MID)
    py += Inches(0.96)

# Right column: the MCP server advantage
rx = Inches(6.98)
rw = Inches(5.85)
R(sl, rx, ly, rw, Inches(0.36), PRIMARY)
T(sl, "Shared MCP Server  (this tool)",
  rx + Inches(0.15), ly + Inches(0.07), rw - Inches(0.3), Inches(0.25),
  size=11, color=WHITE, bold=True)

advantages = [
    ("Shared URL cache", "One fetch serves all concurrent agents. 10 agents needing the same page = 1 render + 9 cache hits. Configurable TTL. Reduces egress cost and latency at scale."),
    ("Tier memoization", "The server learns which domains need Playwright vs HTTP. That knowledge is shared across all agents and persists across sessions — amortised at deployment."),
    ("Single POPIA chokepoint", "All web traffic passes through one auditable service. One place to enforce engine profiles, log requests, rotate provider keys, and respond to legal holds."),
    ("Agent-optimized markdown", "The markdown-for-agents library produces structured, boilerplate-free Markdown — consistent format regardless of source site. Every agent sees the same quality output."),
]
py = ly + Inches(0.45)
for title, body in advantages:
    R(sl, rx, py, rw, Inches(0.88), WHITE)
    R(sl, rx, py, Inches(0.08), Inches(0.88), GREEN)
    T(sl, title, rx + Inches(0.2), py + Inches(0.06), rw - Inches(0.28), Inches(0.3),
      size=12, color=DARK, bold=True)
    T(sl, body,  rx + Inches(0.2), py + Inches(0.38), rw - Inches(0.28), Inches(0.46),
      size=10, color=MID)
    py += Inches(0.96)

# Bottom insight
R(sl, Inches(0.5), ct + Inches(5.35), Inches(12.33), Inches(0.62), PRIMARY_TINT)
R(sl, Inches(0.5), ct + Inches(5.35), Inches(0.1), Inches(0.62), PRIMARY)
T(sl, "Claude Code can also be configured to use this MCP server — replacing its built-in WebFetch "
       "with a Playwright-rendered, cache-backed, POPIA-profiled fetch. Developer and production agent "
       "use the same shared infrastructure.",
  Inches(0.7), ct + Inches(5.42), Inches(12.0), Inches(0.48),
  size=11, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — EMPIRICAL COMPARISON: CLAUDE CODE TOOLS vs THIS MCP SERVER
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "Claude Code Standard Tools vs This MCP Server — Empirical Comparison")

T(sl, "All claims below are verifiable from the source code in this repository "
       "(src/render/ladder.ts, src/render/browserPool.ts, src/fetcher.ts).",
  Inches(0.5), ct, Inches(12.4), Inches(0.38),
  size=12, color=MID, italic=True)

# Two comparison tables side by side
# ── Table 1: fetch / WebFetch ──────────────────────────────────────────────────
T(sl, "Web Fetch",
  Inches(0.5), ct + Inches(0.5), Inches(6.1), Inches(0.35),
  size=14, color=DARK, bold=True)

fetch_rows = [
    ("Capability",                    "Claude Code WebFetch",           "fetch()  via this MCP server"),
    ("Rendering engine",              "HTTP client — no browser engine",  "Three-tier ladder: HTTP → Lightpanda → Playwright/Chromium"),
    ("JavaScript execution",          "No",                              "Yes — auto-escalated when heuristic detects JS-rendered page"),
    ("SPA / React / Angular support", "No — returns pre-render HTML",    "Yes — heuristic detects empty root mount, escalates to Playwright"),
    ("Output format",                 "Raw HTML or minimally processed text", "Agent-optimized Markdown via markdown-for-agents library"),
    ("Shared URL cache",              "No — per-session, no persistence","Yes — LRU cache shared across all agents (configurable TTL)"),
    ("Tier memoization",              "No",                              "Yes — server learns which domains need Playwright"),
    ("Resource blocking",             "No",                              "Yes — blocks images, fonts, ads by default (RENDER_BLOCK_RESOURCES)"),
    ("POPIA engine profile",          "Not configurable",                "Yes — clean / balanced / external profiles"),
    ("Designed for",                  "Single developer session",        "Shared enterprise infrastructure"),
]

col_w1 = [Inches(2.35), Inches(1.82), Inches(2.2)]
tbl1 = sl.shapes.add_table(
    len(fetch_rows), 3,
    Inches(0.5), ct + Inches(0.92),
    sum(col_w1), Inches(4.6)
).table
for ci, w in enumerate(col_w1):
    tbl1.columns[ci].width = w

for ri, row in enumerate(fetch_rows):
    for ci, val in enumerate(row):
        cell = tbl1.cell(ri, ci)
        cell.text = val
        is_header = ri == 0
        if is_header:
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        elif ci == 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = STRIPE if ri % 2 == 0 else WHITE
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF4) if ri % 2 == 0 else WHITE
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(9.5); run.font.name = "Calibri"
        run.font.bold = is_header or ci == 0
        run.font.color.rgb = WHITE if is_header else (MID if ci == 1 and not is_header else DARK)
        p.alignment = PP_ALIGN.LEFT

# ── Table 2: search / WebSearch ──────────────────────────────────────────────────
T(sl, "Web Search",
  Inches(6.98), ct + Inches(0.5), Inches(6.1), Inches(0.35),
  size=14, color=DARK, bold=True)

search_rows = [
    ("Capability",                "Claude Code WebSearch",          "search()  via this MCP server"),
    ("Provider",                  "Anthropic-controlled endpoint",  "Configurable: SearXNG, Brave API, DuckDuckGo"),
    ("Self-hostable provider",    "No",                             "Yes — SearXNG runs on your own infra"),
    ("Data residency",            "Not configurable",               "Fully SA-deployable — no cross-border calls"),
    ("Multi-provider fan-out",    "No — single provider",           "Yes — parallel fan-out with RRF merge"),
    ("Result caching",            "No",                             "Yes — shared cache, configurable TTL"),
    ("POPIA profile",             "Not configurable",               "Three pre-configured profiles"),
    ("Audit log",                 "No — per-session only",          "Yes — all queries through one service"),
    ("Designed for",              "Single developer session",       "Shared enterprise infrastructure"),
]

col_w2 = [Inches(2.1), Inches(1.82), Inches(2.35)]
tbl2 = sl.shapes.add_table(
    len(search_rows), 3,
    Inches(6.98), ct + Inches(0.92),
    sum(col_w2), Inches(4.6)
).table
for ci, w in enumerate(col_w2):
    tbl2.columns[ci].width = w

for ri, row in enumerate(search_rows):
    for ci, val in enumerate(row):
        cell = tbl2.cell(ri, ci)
        cell.text = val
        is_header = ri == 0
        if is_header:
            cell.fill.solid(); cell.fill.fore_color.rgb = PRIMARY
        elif ci == 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = STRIPE if ri % 2 == 0 else WHITE
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF4) if ri % 2 == 0 else WHITE
        p = cell.text_frame.paragraphs[0]
        run = p.runs[0]
        run.font.size = Pt(9.5); run.font.name = "Calibri"
        run.font.bold = is_header or ci == 0
        run.font.color.rgb = WHITE if is_header else (MID if ci == 1 and not is_header else DARK)
        p.alignment = PP_ALIGN.LEFT


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ENTERPRISE KNOWLEDGE INDEX (ROADMAP)
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "The Next Chapter: Enterprise Knowledge Index")

T(sl, "Phase 2 roadmap — the same MCP interface, extended to index your internal document corpus.",
  Inches(0.5), ct, Inches(12.4), Inches(0.4),
  size=14, color=MID)

# Left: concept
lx = Inches(0.5)
lw = Inches(5.9)
ly = ct + Inches(0.55)
T(sl, "What it does:",
  lx, ly, lw, Inches(0.35), size=13, color=DARK, bold=True)
T(sl, "An agent calls search_knowledge(\"annual leave policy\") and receives relevant "
       "chunks from SharePoint and Confluence in milliseconds — without a live API call. "
       "The index is continuously updated by background connectors and filtered to "
       "documents the querying user is permitted to see.",
  lx, ly + Inches(0.38), lw, Inches(1.3), size=12, color=MID)

T(sl, f"Why {BRAND['org']} is uniquely positioned:",
  lx, ly + Inches(1.85), lw, Inches(0.35), size=13, color=DARK, bold=True)
bullets = [
    "OpenShift already running — $0 marginal infra for the knowledge index",
    "M365 already licensed — SharePoint is the highest-value corpus",
    "POPIA requirement makes Glean and Azure AI Search difficult to justify",
    "No SA-deployed, MCP-native, POPIA-compliant knowledge index exists in market",
    "Existing governance pack is the template — minimal new documentation required",
]
for i, b in enumerate(bullets):
    R(sl, lx, ly + Inches(2.27) + i * Inches(0.57), Inches(0.09), Inches(0.09), PRIMARY)
    T(sl, b, lx + Inches(0.22), ly + Inches(2.2) + i * Inches(0.57), lw - Inches(0.3), Inches(0.52),
      size=12, color=MID)

# Right: phase/source tiers
rx = Inches(7.1)
rw = Inches(5.75)
T(sl, "Crawl Sources — Phased Delivery",
  rx, ct + Inches(0.55), rw, Inches(0.38), size=13, color=DARK, bold=True)

tiers = [
    ("Phase 1 — MVP  (6–8 weeks)",
     ["SharePoint / Microsoft 365", "Confluence (Cloud + Data Center)", "Internal intranet portals"],
     GREEN),
    ("Phase 2 — ACL + Vector Search  (8–12 wks)",
     ["Per-user ACL enforcement (Entra ID groups)", "ServiceNow knowledge base", "Semantic hybrid search (BM25 + embeddings)"],
     BLUE),
    ("Phase 3 — Production Hardening",
     ["Jira / GitHub / GitLab wikis", "Data-subject deletion (POPIA §23–25)", "Formal SLOs and admin dashboard"],
     AMBER),
]

ty = ct + Inches(1.03)
for title, items, c in tiers:
    R(sl, rx, ty, rw, Inches(0.36), c)
    T(sl, title, rx + Inches(0.15), ty + Inches(0.05), rw - Inches(0.3), Inches(0.27),
      size=10, color=WHITE, bold=True)
    for j, item in enumerate(items):
        iy = ty + Inches(0.38) + j * Inches(0.38)
        R(sl, rx, iy, rw, Inches(0.37), STRIPE if j % 2 == 0 else WHITE)
        R(sl, rx, iy, Inches(0.07), Inches(0.37), c)
        T(sl, item, rx + Inches(0.18), iy + Inches(0.06), rw - Inches(0.25), Inches(0.27),
          size=11, color=MID)
    ty += Inches(0.38) + len(items) * Inches(0.38) + Inches(0.2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — WHY THIS IS THE RIGHT CHOICE
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, f"Why This Is the Right Choice for {BRAND['org']}")

reasons = [
    ("1", "POPIA by Design",
     f"Self-hosted on {BRAND['org']} SA infrastructure. No Section 72 cross-border transfer. Full data flow documentation, "
     "POPIA assessment, and threat model already in the repository. No legal risk of the kind "
     "Tavily, Firecrawl, or Glean SaaS introduce."),
    ("2", "Runs on Infrastructure You Already Own",
     f"Mode G on existing OpenShift: $0 marginal infrastructure cost. Mode F on ECS: the same runtime "
     f"{BRAND['platform']} already uses. No new procurement, no new cloud accounts, no new approval cycles."),
    ("3", "Audit-Ready, Not Audit-Pending",
     "Complete governance pack: POPIA assessment, data flow, threat model, security findings register, runbook, "
     "SLO template, and cost model. A security or privacy review can start today. "
     "Tavily and Firecrawl ship none of this."),
    ("4", "MCP-Native — Standard Protocol",
     f"MCP is the emerging standard for agent tool integration. {BRAND['platform']}, Claude, and most modern LLM "
     "frameworks support it natively. No bespoke integration code. "
     "Agents call fetch() or search() directly today; search_knowledge() lands in Phase 2."),
    ("5", "A Platform, Not a Point Solution",
     "Web fetch and search today. Enterprise knowledge index (SharePoint + Confluence) next. "
     "The same MCP interface, the same governance pack, the same deployment model. "
     "Each phase extends the platform rather than replacing it."),
]

rh = Inches(1.08)
gap = Inches(0.065)
ry = ct + Inches(0.1)

for num, title, body in reasons:
    R(sl, Inches(0.5),  ry, Inches(0.7), rh, PRIMARY)
    T(sl, num, Inches(0.5), ry + Inches(0.2), Inches(0.7), Inches(0.58),
      size=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    R(sl, Inches(1.22), ry, Inches(11.6), rh, WHITE)
    T(sl, title, Inches(1.38), ry + Inches(0.09), Inches(11.2), Inches(0.36),
      size=13, color=DARK, bold=True)
    T(sl, body,  Inches(1.38), ry + Inches(0.48), Inches(11.2), Inches(0.54),
      size=11, color=MID)
    ry += rh + gap


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — NEXT STEPS
# ═══════════════════════════════════════════════════════════════════════════════
sl = new_slide()
ct = header(sl, "Proposed Next Steps")

steps_data = [
    ("Week 1–2",   "Confirm Deployment",
     ["Confirm Mode G (OpenShift) as target",
      "Initiate Microsoft Graph app registration",
      "Identify first " + BRAND["platform"] + " agent workload"],
     BLUE),
    ("Week 2–4",   "Deploy Phase 1 MVP",
     ["Deploy MCP server to OpenShift",
      "Connect SearXNG and/or Brave API with DPA",
      BRAND["platform"] + " agents using fetch() + search() in staging"],
     GREEN),
    ("Week 4–8",   "Internal Connectors",
     ["SharePoint connector (Graph API, delta queries)",
      "Confluence connector (Cloud or Data Center)",
      "Intranet portal crawl via Playwright pipeline"],
     AMBER),
    ("Month 3–6",  "ACL + Vector Search",
     ["Entra ID per-user ACL enforcement",
      "Vector embeddings + hybrid BM25/semantic search",
      "ServiceNow knowledge base connector"],
     PRIMARY),
]

sw2 = Inches(2.97)
gap2 = Inches(0.14)
sx2  = Inches(0.5)
sy2  = ct + Inches(0.28)
sh2  = Inches(4.85)

for i, (tf, title, bullets, c) in enumerate(steps_data):
    bx = sx2 + i * (sw2 + gap2)
    shp = sl.shapes.add_shape(1, bx, sy2, sw2, sh2)
    shp.fill.solid(); shp.fill.fore_color.rgb = WHITE
    shp.line.color.rgb = BORDER; shp.line.width = Pt(0.5)
    R(sl, bx, sy2, sw2, Inches(0.08), c)
    R(sl, bx + Inches(0.15), sy2 + Inches(0.16), Inches(1.6), Inches(0.28), c)
    T(sl, tf, bx + Inches(0.15), sy2 + Inches(0.16), Inches(1.6), Inches(0.28),
      size=8, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    T(sl, title, bx + Inches(0.15), sy2 + Inches(0.55), sw2 - Inches(0.3), Inches(0.5),
      size=13, color=c, bold=True)
    R(sl, bx + Inches(0.15), sy2 + Inches(1.1), sw2 - Inches(0.3), Inches(0.02), c)
    for bi, b in enumerate(bullets):
        by2 = sy2 + Inches(1.22) + bi * Inches(1.12)
        R(sl, bx + Inches(0.2), by2 + Inches(0.4), Inches(0.09), Inches(0.09), c)
        T(sl, b, bx + Inches(0.4), by2 + Inches(0.1), sw2 - Inches(0.58), Inches(0.9),
          size=11, color=MID)

R(sl, Inches(0.5), sy2 + sh2 + Inches(0.18), Inches(12.33), Inches(0.65), PRIMARY_TINT)
R(sl, Inches(0.5), sy2 + sh2 + Inches(0.18), Inches(0.1), Inches(0.65), PRIMARY)
T(sl, f"Precondition: identify the specific {BRAND['platform']} agent workload that is blocked today without web access. "
       "That workload is the business case. Without a named workload, defer until the use case is concrete.",
  Inches(0.7), sy2 + sh2 + Inches(0.25), Inches(12.0), Inches(0.5),
  size=11, color=DARK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE NOTES — added after all slides are built
# ═══════════════════════════════════════════════════════════════════════════════

SEP = "\n" + "─" * 76 + "\n\n"


def add_notes(slide_index, text):
    notes_slide = prs.slides[slide_index].notes_slide
    notes_slide.notes_text_frame.text = text


# ── Slide 1: Title ────────────────────────────────────────────────────────────
add_notes(0, _b(
    "WHAT TO SAY\n"
    "• Welcome. This presentation introduces a self-hosted web intelligence capability\n"
    "  built for your organisation's agent platform — enabling agents to access live web content and,\n"
    "  on the roadmap, internal documents such as SharePoint and Confluence.\n"
    "• I'll cover what the tool does, why external vendors create POPIA exposure,\n"
    "  what it costs, how it deploys, and what the next steps look like.\n"
    "• Three headline facts set the frame: $0 marginal infrastructure on OpenShift,\n"
    "  MIT open-source licence, and POPIA compliance built into the design.\n"
    "• By the end, I'd like to leave the room with a clear decision on whether to\n"
    "  move this to staging." + SEP +
    "CONTEXT & DETAIL\n"
    "The tool is named markdown-for-agents-mcp. It is a Node.js MCP server that gives\n"
    "AI agents the ability to fetch web pages and run search queries — returning clean,\n"
    "structured Markdown that LLMs consume directly without additional post-processing.\n\n"
    "The '$0 marginal infra' claim applies specifically to the existing OpenShift cluster.\n"
    "OpenShift licensing and the underlying compute are already a sunk cost. This workload\n"
    "adds negligible resource requirements at projected volumes — no new procurement needed.\n\n"
    "The 'POPIA compliant by design' claim refers to architectural choices (self-hosted,\n"
    "SA infrastructure, no cross-border data transfer) and a complete 10-document\n"
    "governance pack already in the repository. The tool has not yet been through a full\n"
    "production authorisation cycle — that step follows this approval decision." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Is this already in production?\n"
    "A: No. The code is production-quality — 1,237 tests passing across 70 test files —\n"
    "   but PRODUCTION_AUTHORISATION.md currently reflects pre-production status.\n"
    "   Completing the authorisation requires a live deployment, SLO measurement,\n"
    "   and formal DPO sign-off.\n\n"
    "Q: Who built this and who owns it?\n"
    "A: Built by the the agent platform team. MIT-licensed, meaning your organisation owns it\n"
    "   completely and can modify or redistribute it without any vendor restriction."
))

# ── Slide 2: The Opportunity ──────────────────────────────────────────────────
add_notes(1, _b(
    "WHAT TO SAY\n"
    "• AI agents are powerful reasoners — but only within the bounds of their training\n"
    "  data. That data has a fixed cutoff, typically 6–12 months before model release.\n"
    "• For questions about current regulation, competitor pricing, or live internal\n"
    "  policies, an agent without web access is reasoning from stale context.\n"
    "• This tool provides three capabilities: fetch any public URL as clean Markdown,\n"
    "  search the web via configurable providers, and — on the roadmap — query a\n"
    "  pre-built index of internal documents.\n"
    "• The third card carries an amber 'Roadmap' badge — it is not a current capability.\n"
    "  Web fetch and search are live today. I'll cover the knowledge index on Slide 12." + SEP +
    "CONTEXT & DETAIL\n"
    "Web Fetch in practice: when an agent calls fetch('https://fsca.co.za/...'), the\n"
    "tool makes an HTTP request, detects whether the page requires JavaScript rendering,\n"
    "converts the rendered HTML to clean Markdown, and returns it. The agent reasons\n"
    "over current document text — not training data.\n\n"
    "The '~60% resource blocking' estimate is based on typical web page composition.\n"
    "Blocking images, fonts, advertising scripts, and third-party tracking reduces\n"
    "download size, NAT egress cost, and rendering latency on every request. Sites with\n"
    "heavy media may see more; simple text pages less.\n\n"
    "Enterprise Knowledge Index: requires Phase 2 development — SharePoint/Confluence\n"
    "connectors, a persistent search index, and per-user ACL enforcement. Estimated\n"
    "build: 6–8 weeks for Phase 1 MVP (broad-access content), plus 8–12 additional\n"
    "weeks for per-user ACL enforcement (the technically complex part)." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: What is the latency for a typical web fetch?\n"
    "A: Static pages via the HTTP tier: 100–500ms. JavaScript-rendered pages via\n"
    "   Playwright: 1–4 seconds. Cached results: under 10ms. The URL cache means\n"
    "   subsequent requests to the same URL are nearly instantaneous.\n\n"
    "Q: Does the tool handle authentication for internal sites?\n"
    "A: Not in Phase 1. For public web, no authentication is needed. For internal\n"
    "   SharePoint and Confluence (Phase 2), authentication is handled at the connector\n"
    "   level using OAuth service principals with the Microsoft Graph API."
))

# ── Slide 3: Why Not External Vendors ─────────────────────────────────────────
add_notes(2, _b(
    "WHAT TO SAY\n"
    "• The POPIA argument is structural, not a preference. Under Section 72, sending\n"
    "  document content to a US-hosted service is a cross-border transfer of personal\n"
    "  information. For a regulated SA operator, this requires a legal justification\n"
    "  that may not be available for all data categories.\n"
    "• The cost argument is secondary but real: Tavily Growth at 10,000 queries/day\n"
    "  costs approximately R330,000 per year. And a workload crossing a Firecrawl tier\n"
    "  triggers a 4× overnight increase with no warning.\n"
    "• Neither Tavily nor Firecrawl ships a POPIA assessment, data flow, or runbook.\n"
    "  Our security review process requires all of these before production approval.\n"
    "• Self-hosting is not the hard option here. The tool is already built and the\n"
    "  governance pack is already written." + SEP +
    "CONTEXT & DETAIL\n"
    "POPIA Section 72 analysis: the key question is whether queries to Tavily or\n"
    "Firecrawl constitute 'processing' of personal information as defined in POPIA.\n"
    "If the queries contain names, identifiers, or references to identified individuals\n"
    "(common in agent workflows where the agent acts on behalf of a named user), then\n"
    "transmitting them to a foreign service triggers Section 72. The Section 72 gateways\n"
    "(binding corporate rules, adequacy finding, consent) may or may not be available\n"
    "depending on the specific data category and the DPO's assessment. Self-hosting\n"
    "eliminates the question entirely.\n\n"
    "Cost calculation: $0.005 × 10,000 queries/day × 365 days = $18,250/year.\n"
    "At R18/$ (approximate) = R328,500. The R330,000 figure in the slide is correct.\n\n"
    "Firecrawl tier jump: Standard (100k pages) = $83/mo → Growth (500k pages)\n"
    "= $333/mo — a 4× increase. A workload that grows from 90k to 110k pages in a\n"
    "single month triggers this jump with no grace period.\n\n"
    "Firecrawl AGPL-3.0 note: AGPL extends the GPL copyleft requirement to network\n"
    "services. Modifying Firecrawl and running it as an internal service requires\n"
    "releasing the modified source code — a legal obligation many enterprise legal\n"
    "teams consider unacceptable." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Can't we just sign Tavily's Data Processing Agreement?\n"
    "A: A properly drafted DPA can satisfy Section 72 via the binding-agreement gateway,\n"
    "   provided it offers 'substantially similar' protection to POPIA. Whether the DPA\n"
    "   Tavily provides meets that standard — and for all data categories your organisation agents\n"
    "   transmit — is the DPO's determination to make, not an engineering decision.\n\n"
    "Q: What about Azure South Africa North?\n"
    "A: Microsoft's SA-region data centres keep storage in SA, but AI model inference\n"
    "   may route through non-SA regions. Self-hosting means no AI processing occurs\n"
    "   externally at all — the retrieval and Markdown conversion happen locally,\n"
    "   and the LLM inference is In-House LLM's responsibility."
))

# ── Slide 4: Introducing the Tool ─────────────────────────────────────────────
add_notes(3, _b(
    "WHAT TO SAY\n"
    "• The tool is a single Node.js process that speaks the Model Context Protocol.\n"
    "  It exposes web fetch, web search, and — in Phase 2 — knowledge index queries.\n"
    "• Seven runtime dependencies. No database, no message queue, no separate Redis.\n"
    "  Compare that to Firecrawl's 6–7 deployed services. Simpler means fewer failure\n"
    "  modes and a smaller attack surface.\n"
    "• MIT licence — your organisation can fork, modify, and redistribute without any AGPL\n"
    "  obligation. We own this tool completely.\n"
    "• 1,237 tests passing. This is not a prototype. It has production-grade test\n"
    "  coverage including integration tests, security tests, and cost-model guards." + SEP +
    "CONTEXT & DETAIL\n"
    "The 7 runtime dependencies are: @modelcontextprotocol/sdk (MCP protocol),\n"
    "markdown-for-agents (HTML-to-Markdown conversion), playwright (Chromium rendering),\n"
    "prom-client (Prometheus metrics), robots-parser (robots.txt compliance), undici\n"
    "(HTTP client), and zod (schema validation). All are MIT/Apache/BSD licensed.\n\n"
    "Playwright is a Microsoft open-source project (Apache 2.0). It drives Chromium for\n"
    "JavaScript rendering. The tool uses Playwright only for the top tier of its render\n"
    "ladder — Chromium is only instantiated when the heuristic determines the page is\n"
    "JavaScript-heavy. Most pages use the cheaper HTTP tier.\n\n"
    "The '<5 min to connect the agent platform' claim refers to adding the MCP server URL to\n"
    "an MCP client's configuration — a one-time step, not ongoing development work.\n\n"
    "Test coverage breakdown: 70 test files, 1,237 tests passed, 51 skipped.\n"
    "Skipped tests require a live browser or specific network conditions. The test\n"
    "suite includes src/costAnalysis.test.ts — a guard that verifies every dollar\n"
    "figure in the documentation traces back to the reproducible cost model." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Is this production-ready today?\n"
    "A: The code is production-quality. The deployment has not yet been stood up —\n"
    "   PRODUCTION_AUTHORISATION.md is pre-production. SLOs have not been measured\n"
    "   (a live deployment is required). The governance pack provides the framework\n"
    "   for a security review; that review has not yet been formally completed.\n\n"
    "Q: Who maintains this going forward?\n"
    "A: The the agent platform team. The MIT licence means no vendor dependency. Ongoing\n"
    "   maintenance is estimated at 0.015–0.05 FTE depending on incident rate and\n"
    "   whether CAB approval is required for each change."
))

# ── Slide 5: How It Works ─────────────────────────────────────────────────────
add_notes(4, _b(
    "WHAT TO SAY\n"
    "• The pipeline is five steps: agent calls a tool, the server picks a provider,\n"
    "  the page is rendered at the appropriate tier, HTML is converted to Markdown,\n"
    "  and the result is returned.\n"
    "• The innovation in step 3 is the render ladder. The server does not always run\n"
    "  Playwright — that would be slow and expensive. It tries a fast HTTP fetch first.\n"
    "  If the response looks like an SPA or requires JavaScript, it escalates.\n"
    "• The server learns which domains need Playwright. That knowledge persists across\n"
    "  all agent sessions — the second request to the same domain skips the escalation\n"
    "  overhead entirely.\n"
    "• The URL cache sits above this whole pipeline. If any agent has fetched the same\n"
    "  URL within the TTL window, the result is returned in under 10ms." + SEP +
    "CONTEXT & DETAIL\n"
    "Render tier details (src/render/ladder.ts):\n"
    "  • HTTP tier: uses undici (Node.js HTTP client). Fetches raw HTML. ~50–200ms.\n"
    "    Used for static pages, APIs, and well-structured sites.\n"
    "  • Lightpanda tier: lightweight browser engine. Executes some JavaScript without\n"
    "    full Chromium overhead. Handles lightly-dynamic pages faster than Playwright.\n"
    "  • Playwright tier: full Chromium. Handles React, Vue, Angular, SharePoint modern\n"
    "    pages, Confluence with macros. Slowest (~1–4s) but handles everything.\n\n"
    "Tier memoization (src/render/ladder.ts): a probabilistic 5% decay probability\n"
    "ensures the server occasionally re-tests memoised tiers, preventing stale\n"
    "assignments after a site is rebuilt. This is an intentional design choice.\n\n"
    "Heuristic escalation (src/render/heuristic.ts): the needsEscalation() function\n"
    "checks for SPA_MARKERS, empty root mounts (<div id='root'></div>), script tag\n"
    "count, bot challenge content, and text-to-HTML ratio before escalating.\n\n"
    "markdown-for-agents library: purpose-built for LLM consumption. Preserves\n"
    "semantic structure (headings, tables, code blocks, lists) while stripping\n"
    "navigation chrome, cookie banners, and other boilerplate." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: What is the URL cache TTL?\n"
    "A: Default SEARCH_CACHE_TTL_MS = 3,600,000ms (1 hour). Configurable per\n"
    "   deployment via environment variable. For rapidly changing content, reduce\n"
    "   the TTL; for stable reference documents, a longer TTL reduces cost.\n\n"
    "Q: What happens when a page is behind a login?\n"
    "A: Phase 1 cannot authenticate to password-protected pages. For internal\n"
    "   SharePoint and Confluence (Phase 2), authentication is handled at the\n"
    "   connector level using OAuth service principals.\n\n"
    "Q: Can sophisticated anti-bot systems block the tool?\n"
    "A: Some behavioural fingerprinting systems may block Playwright. The tool\n"
    "   detects common bot-challenge patterns (Cloudflare, JS challenges) and returns\n"
    "   a clear error rather than the challenge page. This is a documented open finding."
))

# ── Slide 6: POPIA & Data Sovereignty ─────────────────────────────────────────
add_notes(5, _b(
    "WHAT TO SAY\n"
    "• POPIA compliance here is not a checkbox — it's a structural property of the\n"
    "  architecture. Self-hosted on your organisation infrastructure means personal information\n"
    "  never leaves South African borders. Section 72 never triggers.\n"
    "• Even AWS af-south-1 keeps data in SA for storage, but AI model inference may\n"
    "  route through non-SA regions. With this tool, no AI inference occurs externally\n"
    "  at all — that is In-House LLM's responsibility.\n"
    "• The governance pack on the right is the most immediately useful artefact for\n"
    "  a security review. Ten documents covering every dimension a reviewer needs.\n"
    "• One note of transparency: PRODUCTION_AUTHORISATION.md currently reflects\n"
    "  pre-production status. The framework is complete — finishing it requires a live\n"
    "  deployment and formal DPO sign-off, which is the next step after today." + SEP +
    "CONTEXT & DETAIL\n"
    "POPIA Section 19 (Security Safeguards) requires 'appropriate, reasonable technical\n"
    "and organisational measures' to protect personal information. For AI tools:\n"
    "  • Where does the data go? Self-hosted: stays within your organisation's network perimeter.\n"
    "  • Who has access? Only the platform team — no third-party vendor access.\n"
    "  • Is it encrypted in transit? Yes — TLS for all connections.\n"
    "  • Is it logged and auditable? Yes — structured logging, Prometheus metrics.\n"
    "  • What are the retention policies? Configurable cache TTL; no persistent storage\n"
    "    in Phase 1.\n\n"
    "The three engine profiles:\n"
    "  'clean':    No external search providers. Fetches only target URLs the agent\n"
    "              explicitly provides. Outbound requests go to those URLs only.\n"
    "  'balanced': Adds SearXNG — self-hosted within your organisation's infrastructure.\n"
    "              No external search API calls; SearXNG runs on your organisation infra.\n"
    "  'external': Adds Brave Search API (US-hosted). Requires a Section 72 DPA.\n"
    "              Search queries (not document content) go to Brave." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Has the DPO reviewed this?\n"
    "A: Not yet — this presentation is the first step in that process. POPIA_ASSESSMENT.md\n"
    "   provides the structured impact assessment the DPO review will be based on.\n"
    "   The review can begin immediately from that document.\n\n"
    "Q: PRODUCTION_AUTHORISATION.md says pre-production — does that mean we cannot deploy?\n"
    "A: Not a prohibition. The document defines the conditions for approval. Completing\n"
    "   it requires: a live deployment, SLO measurement, security review sign-off, and\n"
    "   DPO approval. It is the governance framework, not a blocking gate."
))

# ── Slide 7: Total Cost of Ownership ──────────────────────────────────────────
add_notes(6, _b(
    "WHAT TO SAY\n"
    "• At 100,000 pages per month, Mode G at $100/mo ongoing is $17/mo more than\n"
    "  Firecrawl Standard ($83/mo). State that clearly — it is the honest comparison.\n"
    "• The cost picture flips above 100k pages: Firecrawl Standard only covers 100k pages;\n"
    "  the next tier (Growth, 500k pages) costs $333/mo. Mode G at $100/mo wins there.\n"
    "• But cost is not the primary argument for Mode G. POPIA compliance, data\n"
    "  sovereignty, and governance readiness are the deciding factors.\n"
    "• The Glean comparison is striking — approximately 250 times cheaper at any volume." + SEP +
    "CONTEXT & DETAIL\n"
    "Methodology: all figures are produced by scripts/cost-model.py — a reproducible\n"
    "Python model with every input constant sourced and dated. The model is committed\n"
    "to the repository. src/costAnalysis.test.ts validates every dollar figure in the\n"
    "documentation against cost-model-output.json, preventing prose from drifting out\n"
    "of sync with the underlying model.\n\n"
    "Mode G cost components:\n"
    "  • Infrastructure: $0 marginal (OpenShift cluster already paid for).\n"
    "  • Engineering (ongoing best case): 0.015 FTE × R1.5M/yr ≈ $100/mo.\n"
    "  • Engineering (3-year setup amortisation included): $527/mo.\n"
    "  Two break-even points:\n"
    "  - Ongoing $100/mo: Mode G wins vs Firecrawl Growth ($333/mo) at any volume\n"
    "    above 100k pages/month — where the Standard tier limit is crossed.\n"
    "  - Amortised $527/mo: break-even vs Firecrawl Scale ($599/mo) at ~1M pages/month.\n"
    "  During the setup amortisation period, Mode G exceeds Firecrawl Standard on\n"
    "  pure cost — POPIA is the deciding factor at that volume.\n\n"
    "Glean basis: ~$50/user/month (buyer-reported enterprise pricing) × 500 users\n"
    "= $25,000/mo = $300,000/year ≈ R5.4M/year. Actual Glean pricing is custom.\n\n"
    "AWS Kendra (af-south-1, GenAI edition): $1,981/mo based on published pricing\n"
    "at 200k documents and 25k searches/day. Does not include connector engineering,\n"
    "embedding costs, or developer time to configure and maintain." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: If Firecrawl is cheaper, why not just use it?\n"
    "A: POPIA Section 72. Firecrawl is US-hosted. See Slide 3 for the full analysis.\n"
    "   If the DPO concludes Section 72 is not triggered by the specific query content,\n"
    "   Firecrawl Standard at $83/mo is a legitimate option for non-sensitive workloads.\n\n"
    "Q: Does the '$0 infra' claim account for OpenShift licensing?\n"
    "A: No — OpenShift licensing is a sunk cost already paid. Marginal cost means the\n"
    "   additional cost attributable to running this workload on the existing cluster.\n"
    "   Licensing reallocation is a finance question, not an engineering cost.\n\n"
    "SOURCES\n"
    "• All figures: scripts/cost-model-output.json (committed to repository)\n"
    "• Firecrawl pricing: firecrawl.dev/pricing, retrieved 2026-08-24\n"
    "• Tavily pricing: tavily.com/pricing, retrieved 2026-08-24\n"
    "• Engineering rate: R1.5M/yr fully-loaded SA senior engineer (cost model assumption)"
))

# ── Slide 8: Deployment ────────────────────────────────────────────────────────
add_notes(7, _b(
    "WHAT TO SAY\n"
    "• Three deployment modes are supported today. Mode G on OpenShift is the\n"
    "  recommendation — $0 marginal cost and the best POPIA posture.\n"
    "• Mode F on ECS is the path of least resistance for initial testing — it is the\n"
    "  same runtime environment the agent platform already uses, right-sized at ~$251/mo.\n"
    "• In-House LLM is shown as a complementary capability, not a competing one.\n"
    "  In-House LLM handles LLM inference; this tool handles web retrieval. MCP is the\n"
    "  clean separation between them.\n"
    "• 'Deployable in days' refers to the MCP server itself — web fetch and search.\n"
    "  The full knowledge index (Phase 2) is a 6–8 week project." + SEP +
    "CONTEXT & DETAIL\n"
    "OpenShift-specific requirements:\n"
    "  • OpenShift Routes provide Kubernetes Ingress via HAProxy.\n"
    "  • SecurityContextConstraints (SCC): Playwright/Chromium requires specific\n"
    "    security capabilities. The recommended approach is the restricted-v2 SCC\n"
    "    with Chromium's --no-sandbox flag as a fallback. This is an open task.\n"
    "  • KEDA (Kubernetes Event-Driven Autoscaling) is available in OpenShift as an\n"
    "    operator. Provides event-driven scaling based on request queue depth.\n\n"
    "Fargate Spot: instances can be interrupted with 2 minutes' notice. For a stateless\n"
    "MCP server with an in-memory URL cache, a single interruption causes seconds of\n"
    "unavailability — the new instance starts quickly. The ~70% cost reduction is the\n"
    "typical Spot discount; actual savings depend on af-south-1 availability.\n\n"
    "HPA ceiling warning: the shipped ECS topology's HPA max (240 vCPU) costs\n"
    "$11,817/mo — a 14.2× multiple over the desired-count baseline. This is an\n"
    "uncapped spend risk if autoscaling is misconfigured. On OpenShift with KEDA,\n"
    "the scaling policy is more controllable and the ceiling is lower." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: What about high availability — is a single OpenShift pod sufficient?\n"
    "A: For development and staging: yes. For production: minimum 2 replicas for\n"
    "   rolling updates. The MCP server is stateless (all state is in the in-memory\n"
    "   URL cache, which is acceptable to lose on restart). 2 replicas with a liveness\n"
    "   probe provides effective HA.\n\n"
    "Q: Does deploying on OpenShift require CAB approval per deployment?\n"
    "A: Depends on your organisation's change management policy. If the OpenShift cluster has\n"
    "   a standard deployment pipeline, adding a new workload may qualify as a standard\n"
    "   change. CAB approval per deployment raises the ongoing engineering cost from\n"
    "   0.015 FTE to 0.03–0.05 FTE due to administrative overhead."
))

# ── Slide 9: Platform Integration ─────────────────────────────────────────────
add_notes(8, _b(
    "WHAT TO SAY\n"
    "• The integration model is simple: the agent platform's agent framework makes MCP tool\n"
    "  calls. Those calls go to this server. The server handles all the complexity —\n"
    "  rendering, caching, provider selection — and returns clean Markdown.\n"
    "• From the agent's perspective, this is just calling a function: fetch('https://...')\n"
    "  or search('latest FSCA regulation'). No awareness of Playwright or SearXNG.\n"
    "• The stack shows the current Phase 1 state at the bottom: SearXNG and Brave API\n"
    "  are live. SharePoint and Confluence are Phase 2.\n"
    "• The AGENT/INFRA labels show the clean separation — the agent platform owns everything\n"
    "  above the MCP server; this tool owns everything below." + SEP +
    "CONTEXT & DETAIL\n"
    "MCP transport options:\n"
    "  • stdio: MCP server runs as a child process. Communication via stdin/stdout.\n"
    "    Suitable for local development and tightly-coupled deployments.\n"
    "  • SSE (Server-Sent Events): MCP server runs as an HTTP service. Agents connect\n"
    "    via long-lived HTTP connections. Recommended for production server deployments.\n"
    "  • Streamable HTTP (MCP spec 2025-11-25): primary transport for cloud-native\n"
    "    deployments. Supported via the @modelcontextprotocol/sdk package.\n\n"
    "Why MCP matters for the agent platform specifically: MCP provides a standard interface\n"
    "that any compliant agent framework can call. If the agent platform switches LLM providers\n"
    "or agent frameworks in the future, MCP-compatible tools migrate without changes.\n\n"
    "Authentication: MCP 2.1 recommends OAuth 2.1 for authenticated tool servers.\n"
    "For internal deployments within your organisation's network perimeter, bearer token\n"
    "authentication is sufficient. The server's auth model should align with the agent platform's\n"
    "existing identity infrastructure." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Is MCP Anthropic-specific? What if your organisation changes LLM providers?\n"
    "A: MCP was initiated by Anthropic but is an open standard (modelcontextprotocol.io).\n"
    "   It has been adopted by OpenAI, Google, Microsoft, and major LLM frameworks\n"
    "   including LangChain, LlamaIndex, and CrewAI. The protocol is LLM-agnostic.\n\n"
    "Q: What latency does the MCP protocol layer add?\n"
    "A: MCP protocol overhead (serialisation/deserialisation) is negligible — under\n"
    "   5ms. The dominant latency is rendering: HTTP tier 100–500ms, Playwright 1–4s.\n"
    "   Cached results: under 10ms regardless of rendering tier."
))

# ── Slide 10: Designed for Agentic Systems ─────────────────────────────────────
add_notes(9, _b(
    "WHAT TO SAY\n"
    "• There is a fundamental architectural difference between tools designed for one\n"
    "  developer in one session and infrastructure designed for many concurrent agents.\n"
    "• Claude Code's built-in tools are excellent for individual developer use — they\n"
    "  are not designed to be shared enterprise infrastructure. That is a design choice\n"
    "  appropriate to their purpose, not a criticism.\n"
    "• The problems on the left all stem from the same root: per-agent tools have no\n"
    "  shared state, no shared learning, and no shared POPIA control point.\n"
    "• The advantages on the right all follow from deploying a shared service: one\n"
    "  URL cache, one render intelligence layer, one audit boundary.\n"
    "• The callout at the bottom matters: Claude Code itself can be configured to use\n"
    "  this MCP server — developers and production agents share the same infrastructure." + SEP +
    "CONTEXT & DETAIL\n"
    "URL cache sharing example in detail: suppose 10 the agent platform agents are simultaneously\n"
    "researching the same regulatory document at fsca.co.za.\n\n"
    "Without a shared cache:\n"
    "  • 10 separate HTTP requests to fsca.co.za\n"
    "  • 10 separate Playwright renders if the page is JavaScript-heavy\n"
    "  • 10× the bandwidth cost and rendering overhead\n\n"
    "With a shared cache (after the first fetch):\n"
    "  • 1 fetch and render (~2s for the first agent)\n"
    "  • 9 cache hits (under 10ms each)\n"
    "  • Approximately 90% reduction in external requests at the same concurrency\n\n"
    "Tier memoization in practice: the server maintains a {domain → tier} map. When\n"
    "the first request to sharepoint.com escalates to Playwright, that escalation is\n"
    "recorded. Every subsequent request skips the lower tiers — no repeated escalation\n"
    "penalty. That knowledge is shared across all 50 concurrent agents.\n\n"
    "Single POPIA chokepoint: having all web traffic pass through one service means\n"
    "one structured access log for audit, one place to respond to a legal hold, and\n"
    "one place to rotate provider API keys." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Couldn't we achieve the same result with a shared Claude Code agent?\n"
    "A: Claude Code is designed for interactive single-session use. It does not expose\n"
    "   a shared service endpoint that multiple the agent platform agents can call concurrently.\n"
    "   Running 50 Claude Code instances independently would 50× the external request\n"
    "   volume and create 50 separate unaudited egress streams.\n\n"
    "Q: Does the shared cache create a privacy risk across agents?\n"
    "A: The cache stores URL content, not query content. If Agent A fetches a public\n"
    "   document, Agent B gets the same cached content — which is appropriate. For\n"
    "   private or authenticated content (Phase 2), the ACL layer ensures per-user\n"
    "   filtering is applied before any result is returned."
))

# ── Slide 11: Empirical Comparison ────────────────────────────────────────────
add_notes(10, _b(
    "WHAT TO SAY\n"
    "• This slide is designed to be read, not presented. Every claim here is verifiable\n"
    "  from the repository source code — file references are in the slide header.\n"
    "• The comparison is not a criticism of Claude Code. It is designed for a different\n"
    "  purpose — individual developer sessions — and excels at that purpose.\n"
    "• Three rows in the fetch table deserve special attention: rendering engine\n"
    "  (no browser vs three-tier ladder), shared URL cache (no vs yes), and POPIA\n"
    "  engine profile (not configurable vs three profiles).\n"
    "• For search: the critical row is data residency. Claude Code's WebSearch routes\n"
    "  through Anthropic-controlled infrastructure that is not configurable for SA\n"
    "  data residency." + SEP +
    "VERIFICATION GUIDE — HOW TO CONFIRM EACH CLAIM\n\n"
    "Three-tier render ladder:\n"
    "  src/render/ladder.ts — RenderLadder class, TIER_INDEX constant, while loop\n"
    "  showing HTTP → Lightpanda → Playwright escalation.\n\n"
    "JavaScript auto-escalation:\n"
    "  src/render/heuristic.ts — needsEscalation() function checks for SPA_MARKERS,\n"
    "  hasEmptyRootMount(), script count, bot challenge content, text/HTML ratio.\n\n"
    "Shared URL cache:\n"
    "  src/fetcher.ts line 50 — 'export const urlCache = new LazyLRUCache<string>(...)'\n"
    "  This is a module-level singleton shared across all requests in the process.\n\n"
    "Resource blocking:\n"
    "  src/render/browserPool.ts line 186 — the resource blocking intercept.\n"
    "  src/config.ts line 89 — RENDER_BLOCK_RESOURCES: z.string().default('true').\n\n"
    "Claude Code WebFetch — why 'no browser engine':\n"
    "  Claude Code is a CLI tool. It does not bundle a browser. Its WebFetch tool\n"
    "  makes HTTP requests via the Node.js HTTP stack. For an SPA that returns\n"
    "  '<div id=\"root\"></div>' without JavaScript execution, the fetch returns that\n"
    "  empty HTML — which is not useful for LLM consumption.\n\n"
    "Claude Code WebSearch — Anthropic-controlled endpoint:\n"
    "  Based on published Claude Code documentation. The backend search provider is\n"
    "  not independently configurable by the deployer." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Could Anthropic add a shared cache or Playwright to Claude Code in the future?\n"
    "A: Possibly. However, even with those additions, Claude Code would still route\n"
    "   all web traffic through Anthropic's infrastructure — which does not resolve\n"
    "   the POPIA data residency requirement for SA enterprise deployments.\n\n"
    "Q: Is the 'Designed for' row fair — Claude Code was never designed for enterprise\n"
    "   infrastructure?\n"
    "A: Correct — it is a statement of intended use, not a criticism. Claude Code is\n"
    "   the right tool for individual developer sessions. A shared MCP server is the\n"
    "   right tool for production agentic infrastructure. The comparison frames the\n"
    "   choice without overstating it."
))

# ── Slide 12: Enterprise Knowledge Index (Roadmap) ────────────────────────────
add_notes(11, _b(
    "WHAT TO SAY\n"
    "• This slide is roadmap, not current capability. Web fetch and search are live\n"
    "  today. The knowledge index requires Phase 2 build work.\n"
    "• The core value proposition: instead of an agent making a live SharePoint API\n"
    "  call every time it needs a document (1–4 seconds), it queries a pre-built index\n"
    "  and gets the answer in under 200ms.\n"
    "• The three-phase delivery maps to increasing complexity: Phase 1 is broad-access\n"
    "  content (no ACL, deployable in 6–8 weeks). Phase 2 adds per-user ACL enforcement\n"
    "  — the technically complex part, 8–12 additional weeks.\n"
    "• your organisation is well-positioned for this because OpenShift runs, M365 is licensed,\n"
    "  POPIA blocks the SaaS alternatives, and the governance pack is already written." + SEP +
    "CONTEXT & DETAIL\n"
    "Why a knowledge index rather than live API calls?\n\n"
    "Live SharePoint API call latency:\n"
    "  • OAuth token refresh: 0–50ms (cached)\n"
    "  • Graph API request: 100–500ms\n"
    "  • Playwright render of SharePoint modern page: 1–4s\n"
    "  • HTML-to-Markdown conversion: 10–50ms\n"
    "  • Total: 1.1–4.6 seconds per document retrieval\n\n"
    "Knowledge index query latency:\n"
    "  • BM25 search (Phase 1): under 50ms\n"
    "  • Hybrid BM25 + vector search (Phase 2): 50–200ms\n"
    "  • Total: under 200ms for most queries\n\n"
    "For a workflow requiring 5–10 document lookups, this is the difference between\n"
    "5–45 seconds and under 2 seconds of total retrieval time.\n\n"
    "ACL enforcement approach (Phase 2):\n"
    "  Crawl-time ACL snapshot + query-time token enforcement:\n"
    "  1. At crawl time: extract Entra ID groups permitted to access each document\n"
    "     from the Microsoft Graph permissions API.\n"
    "  2. Store as: acl_allowed: [groupId1, groupId2, user@corp.com]\n"
    "  3. At query time: resolve the authenticated user's group memberships\n"
    "     (GET /users/{id}/transitiveMemberOf), cached with 5-minute TTL.\n"
    "  4. Inject an ACL filter into every search — never return content the\n"
    "     querying user is not permitted to see in the source system.\n"
    "  Staleness window: ACL changes are not reflected until the next sync\n"
    "  (1-hour default). This is the same trade-off Glean and Azure AI Search make.\n\n"
    "Market gap: 'No SA-deployed, MCP-native, POPIA-compliant knowledge index exists\n"
    "in market' — based on research conducted 2026-08-25. Glean Customer Hosted requires\n"
    "Glean engineers to retain access. Azure AI Search MCP integration is cloud-only.\n"
    "Atlassian Rovo MCP server is cloud-only (Data Center not supported). Community MCP\n"
    "servers for SharePoint use app-level credentials with no per-user ACL enforcement." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Does your organisation have M365 licensing?\n"
    "A: This slide assumes yes — verify before Phase 2 planning. If your organisation uses a\n"
    "   non-Microsoft productivity suite, the highest-value connector would shift\n"
    "   to Confluence or another internal document platform.\n\n"
    "Q: How does the index stay fresh?\n"
    "A: Microsoft Graph delta queries return only items changed since the last sync\n"
    "   token — efficient for large tenants. Graph webhook subscriptions can trigger\n"
    "   near-real-time re-crawl on document change. Maximum staleness window: 1 hour."
))

# ── Slide 13: Why This Is the Right Choice ─────────────────────────────────────
add_notes(12, _b(
    "WHAT TO SAY\n"
    "• Five reasons, in order of importance. The order matters.\n"
    "• POPIA first: this is the structural argument. It does not matter how good or\n"
    "  cheap the alternatives are if they cannot survive a Section 72 assessment.\n"
    "  Self-hosting on SA infrastructure is the only option that eliminates this\n"
    "  question entirely.\n"
    "• Infrastructure second: $0 marginal cost on OpenShift is a genuine advantage\n"
    "  that compounds over time and with volume.\n"
    "• Governance third: the governance pack is a real differentiator — a security\n"
    "  review that would take weeks with a SaaS vendor can start immediately.\n"
    "• MCP fourth: an open standard, not a proprietary integration. Platform fifth:\n"
    "  each phase extends the same foundation rather than replacing it." + SEP +
    "CONTEXT & DETAIL\n"
    "The 'no legal risk' language is carefully qualified: 'no legal risk of the kind\n"
    "Tavily, Firecrawl, or Glean SaaS introduce.' This means no Section 72 cross-border\n"
    "transfer risk from those specific vendors. The tool itself carries its own risk:\n"
    "  • Open security findings documented in SECURITY_FINDINGS_REGISTER.md\n"
    "  • Pre-production status — not yet through a full formal security review\n"
    "  • Playwright's Chromium component has its own vulnerability surface;\n"
    "    Chromium updates must be tracked\n"
    "  • All 7 runtime deps are MIT/Apache/BSD — no AGPL — but should be re-verified\n"
    "    if new dependencies are added\n\n"
    "The 'Audit-Ready, Not Audit-Pending' claim: the governance pack provides the\n"
    "documentation framework. The production authorisation gate is currently\n"
    "pre-production — meaning the documentation for the review exists and is accurate;\n"
    "it does not mean the review has been completed.\n\n"
    "Why MCP durability matters: vendor-specific integration code becomes a maintenance\n"
    "liability if the platform changes. MCP, as an open standard with broad adoption\n"
    "across Anthropic, OpenAI, Google, and Microsoft, is more likely to survive\n"
    "platform changes than a bespoke API integration." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: What is the total cost of ownership in rand over 3 years?\n"
    "A: Mode G (0.015 FTE, 3-year setup amortisation):\n"
    "   Infrastructure: R0 marginal.\n"
    "   Engineering (3-year amortisation period): $527/mo × 36 = ~$18,972 ≈ R341,500.\n"
    "   Ongoing after amortisation: $100/mo × remaining months.\n"
    "   Firecrawl Standard (100k pages/mo) for reference: $83/mo × 36 = $2,988 ≈ R54k.\n"
    "   At 100k pages/month, Firecrawl is cheaper over 3 years by ~R287,000.\n"
    "   The Mode G case requires either: (a) volume above 1M pages/month, (b) POPIA\n"
    "   requirement that rules out Firecrawl, or (c) valuing data sovereignty and\n"
    "   governance readiness at more than R287,000 over the period.\n\n"
    "Q: Does the existing OpenShift team know about this?\n"
    "A: Not yet — this presentation is the first formal proposal. The next step is\n"
    "   aligning with the OpenShift platform team on namespace, quotas, and SCC."
))

# ── Slide 14: Next Steps ───────────────────────────────────────────────────────
add_notes(13, _b(
    "WHAT TO SAY\n"
    "• Four time-boxed steps. The first two are for this room to decide; the last two\n"
    "  follow automatically if the first two are approved.\n"
    "• Week 1–2 has three items: confirm OpenShift as the deployment target, initiate\n"
    "  the Microsoft Graph app registration, and — most importantly — identify the\n"
    "  specific the agent platform agent workload this will serve first.\n"
    "• The third item is the most important on the entire slide. Without a named\n"
    "  workload and a team that will integrate it, there is no concrete business case.\n"
    "  Building infrastructure ahead of demand is the most common failure mode in\n"
    "  enterprise AI projects.\n"
    "• The red callout is the precondition. If we cannot name the blocked workload\n"
    "  today, the right action is to defer until we can." + SEP +
    "CONTEXT & DETAIL\n"
    "Microsoft Graph app registration (Week 1–2): requires admin consent from the\n"
    "M365 tenant administrator for these application permissions:\n"
    "  • Sites.Read.All — read SharePoint sites and pages\n"
    "  • Files.Read.All — read document library files\n"
    "These require a Global Administrator or Application Administrator to grant consent\n"
    "in the Azure portal. This is a governance step, not a technical one. It should be\n"
    "initiated immediately — it may take days to weeks to obtain approval depending on\n"
    "your organisation's change management process.\n\n"
    "What makes a good first agent workload:\n"
    "  1. The agent currently gives incorrect or outdated answers because it cannot\n"
    "     access a specific external source.\n"
    "  2. The team running the agent has agreed to integrate fetch() or search()\n"
    "     when available.\n"
    "  3. The workload volume is estimable (needed for cost modelling).\n"
    "  4. The content sources are public or broadly accessible (avoids ACL complexity\n"
    "     in Phase 1).\n\n"
    "Example candidate workloads:\n"
    "  • An agent answering questions about FSCA/SARB guidance that cannot access\n"
    "    recent publications due to its training cutoff.\n"
    "  • An internal helpdesk agent that needs to retrieve from the ServiceNow\n"
    "    knowledge base.\n"
    "  • A market intelligence agent that must currently rely on training data\n"
    "    for competitor research." + SEP +
    "ANTICIPATED QUESTIONS\n\n"
    "Q: Who needs to be in the room for Week 1–2 decisions?\n"
    "A: (1) the agent platform team lead — confirm deployment target. (2) M365/Azure AD\n"
    "   administrator — initiate the Graph app registration. (3) Product owner of the\n"
    "   first agent workload — confirm integration intent. DPO should be notified for\n"
    "   awareness; formal DPO sign-off is needed for Phase 2, not Phase 1.\n\n"
    "Q: Is there a demo we can show today?\n"
    "A: The MCP server can be run locally in minutes: npm install && npm run dev.\n"
    "   Any MCP-compatible client (including Claude Code) can then connect and call\n"
    "   fetch() or search() against real URLs. A live demo is straightforward to\n"
    "   set up for a follow-up session.\n\n"
    "Q: What if the OpenShift cluster is at capacity?\n"
    "A: Phase 1 resource requirements are modest: ~0.5 CPU / 1GB RAM at idle, burst\n"
    "   to 2 CPU / 4GB RAM during Playwright rendering. If the cluster is genuinely\n"
    "   at capacity, Mode F (ECS Fargate) is the fallback at ~$251/mo infrastructure."
))


# ═══════════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════════
ROOT = Path(__file__).parent.parent
out_dir = ROOT / "docs" / "marketing"
out_dir.mkdir(parents=True, exist_ok=True)
out_path = out_dir / (BRAND["out_stem"] + ".pptx")
prs.save(str(out_path))
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
